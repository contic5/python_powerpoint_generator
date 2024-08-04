import os
from dotenv import load_dotenv
from pathlib import Path

dotenv_path = Path('environment_variables.env')
load_dotenv(dotenv_path=dotenv_path)

import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from collections import defaultdict
from openai import OpenAI

import google.generativeai as genai
from api_settings import safety_settings

secretkey=os.getenv("secretkey")
api_key=os.getenv("api_key")

client = OpenAI(api_key=secretkey)
genai.configure(api_key=api_key)

#Ais can either be chatgpt or gemini
ais=[]
#ais.append("chatgpt")
ais.append("gemini")

question_file_name="goal_setting"
def add_bullet_slide(prs,question_dict):
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = question_dict["Question"]

    answer=question_dict["Answer"]
    sentences=answer.split(".")
    for i in range(len(sentences)):
        sentence=sentences[i]
        sentence=sentence.strip()
        if i==0:
            tf = body_shape.text_frame
            tf.text = sentence
        else:
            p = tf.add_paragraph()
            p.text = sentence

    p = tf.add_paragraph()
    p.text = question_dict["Source"]
    p.level = 1

def add_questions_slide(prs):
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Questions?"

def add_title_slide(prs,ai):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = question_file_name.replace("_"," ").title()
    subtitle.text = f"By CJ Conti (Generated with Python and {ai})"

def get_answer_from_chatgpt(question):
    system_prompt="""
    You are a Powerpoint presentation assistant, skilled at explaining soft skills concepts. 
    Explain concepts in 3 or fewer short sentences. Do not use periods except at the end of sentences."
    """

    completion = client.chat.completions.create(
    model="gpt-3.5-turbo",
    messages=[
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": question}
    ],
    )

    #print(completion.usage,end="\n\n")
    #print(completion.choices[0].message.content)
    response=completion.choices[0].message.content
    return response

def get_answer_from_gemini(question):
    system_prompt="""
    You are a Powerpoint presentation assistant, skilled at explaining soft skills concepts. 
    Explain concepts in 3 or fewer short sentences. Do not use periods except at the end of sentences."
    """
    model = genai.GenerativeModel('gemini-1.5-flash',safety_settings=safety_settings)
    chat = model.start_chat(history=[])

    response = chat.send_message(system_prompt,safety_settings=safety_settings)

    response = chat.send_message(question,safety_settings=safety_settings)
    return response.text

def blank():
    return ""

def main():
    for ai in ais:
        print(f"Generating presentation on {question_file_name} with {ai}")
        prs=Presentation()
        add_title_slide(prs,ai)

        with open(f"questions/{question_file_name}.txt") as question_file:
            lines=question_file
            for line in lines:
                line=line.strip()
                line=line.replace("\n","")

                question_dict=defaultdict(blank)
                question_dict["Question"]=line

                if ai=="chatgpt":
                    question_dict["Answer"]=get_answer_from_chatgpt(line)
                else:
                    question_dict["Answer"]=get_answer_from_gemini(line)

                print(question_dict["Question"])
                print(question_dict["Answer"])
                print()
                add_bullet_slide(prs,question_dict)
            
        add_questions_slide(prs)
        prs.save(f'presentations/{question_file_name}_{ai}.pptx')

if __name__=="__main__":
    main()